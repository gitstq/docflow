"""PowerPoint to Markdown converter."""

from pathlib import Path

from docflow.core.converters.base import BaseConverter
from docflow.core.models import DocumentMetadata

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PowerPointConverter(BaseConverter):
    """Convert PowerPoint presentations to Markdown."""

    def convert(self, source_path: Path, output_path: Path):
        """Convert a PowerPoint presentation to Markdown."""
        if not HAS_PPTX:
            return self._create_error_result(
                source_path,
                output_path,
                ["python-pptx not installed. Run: pip install python-pptx"],
            )

        try:
            prs = Presentation(source_path)
            markdown_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                markdown_parts.append(f"\n\n---\n\n## Slide {slide_num}\n")

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            markdown_parts.append(f"\n{text}\n")

                    if shape.has_table:
                        md_table = self._format_table(shape.table)
                        markdown_parts.append(f"\n{md_table}\n")

            markdown_content = "".join(markdown_parts).strip()
            metadata = self._extract_metadata(source_path, prs)

            return self._create_success_result(
                source_path,
                output_path,
                markdown_content,
                metadata=metadata,
                pages_processed=len(prs.slides),
            )

        except Exception as e:
            return self._create_error_result(
                source_path, output_path, [f"PowerPoint conversion error: {str(e)}"]
            )

    def _format_table(self, table) -> str:
        """Convert a PowerPoint table to Markdown."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        header = rows[0]
        separator = ["---"] * len(header)
        body = rows[1:]

        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(separator) + " |",
        ]
        for row in body:
            lines.append("| " + " | ".join(row) + " |")

        return "\n".join(lines)

    def _extract_metadata(self, source_path: Path, prs) -> DocumentMetadata:
        """Extract metadata from PowerPoint."""
        props = prs.core_properties

        return DocumentMetadata(
            title=props.title,
            author=props.author,
            subject=props.subject,
            keywords=props.keywords.split(", ") if props.keywords else [],
            creation_date=props.created,
            modification_date=props.modified,
            page_count=len(prs.slides),
            file_size=source_path.stat().st_size if source_path.exists() else None,
        )
